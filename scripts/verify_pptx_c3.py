#!/usr/bin/env python3
"""Independent verification of the C3-generated PPTX via python-pptx.

Reads test-output/allegro-C3.pptx and reports, from the FILE itself (not the
TS in-memory objects): slide count, slide dimensions, and a sample of text /
shape / image elements from a few slides. Serves as the cross-engine check
that the .pptx actually contains what the IR deck describes.
"""
import sys
from pptx import Presentation
import sys

# Force UTF-8 output regardless of the Windows console codepage (cp1250).
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

PPTX = sys.argv[1] if len(sys.argv) > 1 else "test-output/allegro-C3.pptx"

prs = Presentation(PPTX)

errors = []
def check(name, ok, detail=""):
    print(("  [PASS] " if ok else "  [FAIL] ") + name + ((" â€” " + str(detail)) if detail else ""))
    if not ok:
        errors.append(name)

print(f"=== python-pptx verification of {PPTX} ===")

n_slides = len(prs.slides)
check("27 slides", n_slides == 27, f"{n_slides}")

# Slide dimensions (the layout is A4 portrait: 595.3x841.9 pt).
w_in = prs.slide_width / 914400  # EMU -> inches
h_in = prs.slide_height / 914400
w_pt = w_in * 72
h_pt = h_in * 72
check("width â‰ 595.3 pt (8.27in)", abs(w_pt - 595.3) < 2.0, f"{w_pt:.2f}pt/{w_in:.3f}in")
check("height â‰ 841.9 pt (11.69in)", abs(h_pt - 841.9) < 2.0, f"{h_pt:.2f}pt/{h_in:.3f}in")
check("portrait", h_pt > w_pt, f"{w_pt:.1f}x{h_pt:.1f}")

def elem_summary(shape):
    t = shape.text_frame.text if shape.has_text_frame else ""
    kind = "text"
    if shape.shape_type is not None and "AUTO_SHAPE" in str(shape.shape_type):
        kind = "shape"
    if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
        kind = "image"
    x = shape.left / 914400 * 72
    y = shape.top / 914400 * 72
    w = shape.width / 914400 * 72
    h = shape.height / 914400 * 72
    return f"{kind} @({x:.1f},{y:.1f},{w:.1f}x{h:.1f}) " + (t[:60] if t else "")

# Sample a few slides: index 0 (title), 1 (body paragraph), 3 (another), 26 (last).
for idx in (0, 1, 3, 26):
    slide = prs.slides[idx]
    shapes = list(slide.shapes)
    print(f"  --- slide {idx+1}: {len(shapes)} shapes ---")
    for sh in shapes[:6]:
        print(f"      {elem_summary(sh)}")

print("")
if errors:
    print(f"FAILURES PRESENT ({len(errors)})")
    sys.exit(1)
print("ALL PASS")
