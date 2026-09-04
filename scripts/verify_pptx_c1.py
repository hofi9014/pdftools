import sys
from pptx import Presentation
from pptx.util import Emu

path = sys.argv[1] if len(sys.argv) > 1 else "test-output/allegro-C1.pptx"

prs = Presentation(path)

# --- 1) Slide dimensions (independent of the writer) ---
sw = prs.slide_width  # EMU
sh = prs.slide_height
sw_in = sw / 914400.0
sh_in = sh / 914400.0
print("SLIDE  width  EMU=%d  in=%.4f" % (sw, sw_in))
print("SLIDE  height EMU=%d  in=%.4f" % (sh, sh_in))
exp_w, exp_h = 595.3/72.0, 841.9/72.0
print("EXPECT        in=%.4f x %.4f" % (exp_w, exp_h))
assert abs(sw_in - exp_w) < 0.02, "slide width mismatch"
assert abs(sh_in - exp_h) < 0.02, "slide height mismatch"
print("  -> slide size MATCH (A4 portrait)")

# --- 2) Shapes on slide 1: count + geometry + text ---
slide = prs.slides[0]
print("\n# shapes on slide 1:", len(slide.shapes))
for i, sp in enumerate(slide.shapes):
    x_in = sp.left / 914400.0
    y_in = sp.top / 914400.0
    w_in = sp.width / 914400.0
    h_in = sp.height / 914400.0
    txt = sp.text if sp.has_text_frame else ""
    txt = (txt.replace("\n", "\\n")[:60]) if txt else "(none)"
    print("  shape[%d] type=%s x=%.3f y=%.3f w=%.3f h=%.3f text=%r" % (
        i, sp.shape_type, x_in, y_in, w_in, h_in, txt))

assert len(slide.shapes) == 3, "expected exactly 3 elements (title, body, rect) on slide 1"

# text checks
texts = [sp.text for sp in slide.shapes if sp.has_text_frame and sp.text.strip()]
assert any("12 rzeczy" in t for t in texts), "title textbox missing"
assert any("audytu" in t for t in texts), "body textbox missing"
print("\n  -> textboxes present:", texts)

print("\nINDEPENDENT VERIFICATION OK (python-pptx)")